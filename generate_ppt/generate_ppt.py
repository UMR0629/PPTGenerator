#from lxml.xmlerror import text_size
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image


class Generate_ppt:
    def __init__(self, ppt_path):
        self.ppt_path = ppt_path
        self.prs_old = Presentation(ppt_path)
        self.prs_new = self.prs_old
        self.slide_layout = self.prs_new.slide_layouts

    # 去除标题序号
    def process_title(self, text_to_process):
        return re.sub(r'^\d+[\.\s]+', '', text_to_process)

    # 添加封面页
    def add_cover(self, title_text, author_text, date_text):
        slide = self.prs_new.slides.add_slide(self.slide_layout[0])
        title_len = len(title_text)
        print(f"标题长度：{title_len}")
        if title_len < 50:
            text_size = 54
        elif 50 <= title_len < 100:
            text_size = 32
        else:
            text_size = 24
        print(text_size)

        slide.placeholders[0].text = title_text

        for paragraph in slide.placeholders[0].text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(text_size)


        try:
            slide.placeholders[11].text = author_text
            slide.placeholders[10].text = date_text
        except IndexError:
            print("占位符 10 或 11 不存在，跳过赋值")

    # 添加目录页
    def add_menu(self, image_path, menu_items_number, menu_items):
        if menu_items_number == 3:
            self.add_menu_3(image_path, menu_items)
        elif menu_items_number == 4:
            self.add_menu_4(image_path, menu_items)
        elif menu_items_number == 5:
            self.add_menu_5(image_path, menu_items)
        elif menu_items_number == 6:
            self.add_menu_6(image_path, menu_items)
        else:
            print("目录项数量错误，目前仅支持3-6个目录项")

    # 添加目录页（3个目录项）
    def add_menu_3(self, image_path, menu_items):
        slide = self.prs_new.slides.add_slide(self.slide_layout[6])  # 选择布局索引6
        placeholder = slide.placeholders[10]  # 获取占位符

        placeholder.insert_picture(image_path)  # 插入图片

        # 添加目录项
        if not menu_items:
            # 如果目录项为空，则使用默认目录项
            menu_items = ["演示文稿目录1", "演示文稿目录2", "演示文稿目录3"]
        else:
            # 如果目录项超过4个，截取前4个
            if len(menu_items) > 3:
                print("目录项超过3个，只取前3个")
                menu_items = menu_items[:3]
            elif len(menu_items) < 3:
                print("目录项不足3个，使用默认目录项")
                menu_items = ["演示文稿目录1", "演示文稿目录2", "演示文稿目录3"]

        # 遍历目录项，并为每个目录项创建文本框
        for i, text in enumerate(menu_items):
            # 设置文本框的位置和大小
            left, top, width, height = Inches(7.2), Inches(2 + i * 1.2), Inches(5), Inches(0.8)

            # 添加一个文本框到幻灯片中，设置位置和大小
            box = slide.shapes.add_textbox(left, top, width, height)

            # 设置文本框背景为白色
            box.fill.solid()  # 填充颜色
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 背景色为白色

            # 获取文本框的文本框架（TextFrame），用于设置文本属性
            text_frame = box.text_frame
            text_frame.clear()  # 清空任何已有的文本内容

            # 设置文本框中的垂直对齐方式为居中
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # 垂直居中

            # 获取第一个段落，设置该段落的对齐方式为水平居中
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER  # 水平居中

            # 设置该段落的文字样式
            run = p.add_run()  # 为段落添加一个文本运行
            run.text = self.process_title(text)  # 设置文本为目录项的名称
            run.font.size = Pt(24)  # 设置字体大小为24磅
            run.font.bold = True  # 设置字体为粗体
            run.font.color.rgb = RGBColor(68, 84, 106)  # 设置字体颜色为黑色

            # 添加编号框（左侧显示编号）
            num_box = slide.shapes.add_textbox(left - Inches(1), top, Inches(0.8), height)  # 设置编号框的位置和大小
            num_text_frame = num_box.text_frame  # 获取编号框的文本框架
            num_text_frame.clear()  # 清空任何已有的文本内容

            # 设置编号框中的垂直对齐方式为居中
            num_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 设置编号框中的文本段落对齐方式为居中
            p = num_text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER  # 水平居中

            # 设置该段落的编号文本
            run = p.add_run()
            run.text = str(i + 1)  # 编号从1开始
            run.font.size = Pt(28)  # 设置编号字体大小为28磅
            run.font.bold = True  # 设置编号为粗体
            run.font.color.rgb = RGBColor(255, 255, 255)  # 设置编号的字体颜色为白色

            # 设置编号框背景颜色为深色
            num_box.fill.solid()  # 填充编号框颜色
            num_box.fill.fore_color.rgb = RGBColor(54, 69, 79)  # 设置背景色为深色（灰蓝色）

            # 设置编号框的边框颜色为透明
            num_box.line.fill.background()  # 设置编号框边框为透明

    # 添加目录页（4个目录项）
    def add_menu_4(self, image_path, menu_items):
        slide = self.prs_new.slides.add_slide(self.slide_layout[6])  # 选择布局索引6
        placeholder = slide.placeholders[10]  # 获取占位符

        placeholder.insert_picture(image_path)  # 插入图片

        # 添加目录项
        if not menu_items:
            # 如果目录项为空，则使用默认目录项
            menu_items = ["演示文稿目录1", "演示文稿目录2", "演示文稿目录3", "演示文稿目录4"]
        else:
            # 如果目录项超过4个，截取前4个
            if len(menu_items) > 4:
                print("目录项超过4个，只取前4个")
                menu_items = menu_items[:4]
            elif len(menu_items) < 4:
                print("目录项不足4个，使用默认目录项")
                menu_items = ["演示文稿目录1", "演示文稿目录2", "演示文稿目录3", "演示文稿目录4"]

        # 遍历目录项，并为每个目录项创建文本框
        for i, text in enumerate(menu_items):
            # 设置文本框的位置和大小
            left, top, width, height = Inches(7.2), Inches(1.5 + i * 1.2), Inches(5), Inches(0.8)

            # 添加一个文本框到幻灯片中，设置位置和大小
            box = slide.shapes.add_textbox(left, top, width, height)

            # 设置文本框背景为白色
            box.fill.solid()  # 填充颜色
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 背景色为白色

            # 获取文本框的文本框架（TextFrame），用于设置文本属性
            text_frame = box.text_frame
            text_frame.clear()  # 清空任何已有的文本内容

            # 设置文本框中的垂直对齐方式为居中
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # 垂直居中

            # 获取第一个段落，设置该段落的对齐方式为水平居中
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER  # 水平居中

            # 设置该段落的文字样式
            run = p.add_run()  # 为段落添加一个文本运行
            run.text = self.process_title(text)  # 设置文本为目录项的名称
            run.font.size = Pt(24)  # 设置字体大小为24磅
            run.font.bold = True  # 设置字体为粗体
            run.font.color.rgb = RGBColor(68, 84, 106)  # 设置字体颜色为黑色

            # 添加编号框（左侧显示编号）
            num_box = slide.shapes.add_textbox(left - Inches(1), top, Inches(0.8), height)  # 设置编号框的位置和大小
            num_text_frame = num_box.text_frame  # 获取编号框的文本框架
            num_text_frame.clear()  # 清空任何已有的文本内容

            # 设置编号框中的垂直对齐方式为居中
            num_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 设置编号框中的文本段落对齐方式为居中
            p = num_text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER  # 水平居中

            # 设置该段落的编号文本
            run = p.add_run()
            run.text = str(i + 1)  # 编号从1开始
            run.font.size = Pt(28)  # 设置编号字体大小为28磅
            run.font.bold = True  # 设置编号为粗体
            run.font.color.rgb = RGBColor(255, 255, 255)  # 设置编号的字体颜色为白色

            # 设置编号框背景颜色为深色
            num_box.fill.solid()  # 填充编号框颜色
            num_box.fill.fore_color.rgb = RGBColor(54, 69, 79)  # 设置背景色为深色（灰蓝色）

            # 设置编号框的边框颜色为透明
            num_box.line.fill.background()  # 设置编号框边框为透明

    # 添加目录页（5个目录项）
    def add_menu_5(self, image_path, menu_items):
        slide = self.prs_new.slides.add_slide(self.slide_layout[6])  # 选择布局索引6
        placeholder = slide.placeholders[10]  # 获取占位符

        placeholder.insert_picture(image_path)  # 插入图片

        # 添加目录项
        if not menu_items:
            # 如果目录项为空，则使用默认目录项
            menu_items = ["演示文稿目录1", "演示文稿目录2", "演示文稿目录3", "演示文稿目录4", "演示文稿目录5"]
        else:
            # 如果目录项超过5个，截取前5个
            if len(menu_items) > 5:
                print("目录项超过5个，只取前5个")
                menu_items = menu_items[:5]
            elif len(menu_items) < 5:
                print("目录项不足5个，使用默认目录项")
                menu_items = ["演示文稿目录1", "演示文稿目录2", "演示文稿目录3", "演示文稿目录4", "演示文稿目录5"]

        # 遍历目录项，并为每个目录项创建文本框
        for i, text in enumerate(menu_items):
            # 设置文本框的位置和大小
            left, top, width, height = Inches(7.2), Inches(1 + i * 1.2), Inches(5), Inches(0.8)

            # 添加一个文本框到幻灯片中，设置位置和大小
            box = slide.shapes.add_textbox(left, top, width, height)

            # 设置文本框背景为白色
            box.fill.solid()  # 填充颜色
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 背景色为白色

            # 获取文本框的文本框架（TextFrame），用于设置文本属性
            text_frame = box.text_frame
            text_frame.clear()  # 清空任何已有的文本内容

            # 设置文本框中的垂直对齐方式为居中
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # 垂直居中

            # 获取第一个段落，设置该段落的对齐方式为水平居中
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER  # 水平居中

            # 设置该段落的文字样式
            run = p.add_run()  # 为段落添加一个文本运行
            run.text = self.process_title(text)   # 设置文本为目录项的名称
            run.font.size = Pt(24)  # 设置字体大小为24磅
            run.font.bold = True  # 设置字体为粗体
            run.font.color.rgb = RGBColor(68, 84, 106)  # 设置字体颜色为黑色

            # 添加编号框（左侧显示编号）
            num_box = slide.shapes.add_textbox(left - Inches(1), top, Inches(0.8), height)  # 设置编号框的位置和大小
            num_text_frame = num_box.text_frame  # 获取编号框的文本框架
            num_text_frame.clear()  # 清空任何已有的文本内容

            # 设置编号框中的垂直对齐方式为居中
            num_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 设置编号框中的文本段落对齐方式为居中
            p = num_text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER  # 水平居中

            # 设置该段落的编号文本
            run = p.add_run()
            run.text = str(i + 1)  # 编号从1开始
            run.font.size = Pt(28)  # 设置编号字体大小为28磅
            run.font.bold = True  # 设置编号为粗体
            run.font.color.rgb = RGBColor(255, 255, 255)  # 设置编号的字体颜色为白色

            # 设置编号框背景颜色为深色
            num_box.fill.solid()  # 填充编号框颜色
            num_box.fill.fore_color.rgb = RGBColor(54, 69, 79)  # 设置背景色为深色（灰蓝色）

            # 设置编号框的边框颜色为透明
            num_box.line.fill.background()  # 设置编号框边框为透明

    # 添加目录页（6个目录项）
    def add_menu_6(self, image_path, menu_items):
        slide = self.prs_new.slides.add_slide(self.slide_layout[7])
        placeholder = slide.placeholders[10]

        placeholder.insert_picture(image_path)

        if not menu_items:
            menu_items = ["目录1", "目录2", "目录3", "目录4", "目录5", "目录6"]
        else:
            if len(menu_items) > 6:
                print("目录项超过6个，只取前6个")
                menu_items = menu_items[:6]
            elif len(menu_items) < 6:
                print("目录项不足6个，使用默认目录项")
                menu_items = ["目录1", "目录2", "目录3", "目录4", "目录5", "目录6"]

        # 遍历目录项，并为每个目录项创建文本框
        for i, text in enumerate(menu_items):
            # 设置文本框的位置和大小
            if i < 3:
                left, top, width, height = Inches(1.4), Inches(3.2 + i * 1.2), Inches(5), Inches(0.8)
            else:
                left, top, width, height = Inches(8.0), Inches(3.2 + (i - 3) * 1.2), Inches(5), Inches(0.8)

            # 添加一个文本框到幻灯片中，设置位置和大小
            box = slide.shapes.add_textbox(left, top, width, height)

            # 设置文本框背景为白色
            box.fill.solid()  # 填充颜色
            box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 背景色为白色

            # 获取文本框的文本框架（TextFrame），用于设置文本属性
            text_frame = box.text_frame
            text_frame.clear()  # 清空任何已有的文本内容

            # 设置文本框中的垂直对齐方式为居中
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # 垂直居中

            # 获取第一个段落，设置该段落的对齐方式为水平居中
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER  # 水平居中

            # 设置该段落的文字样式
            run = p.add_run()  # 为段落添加一个文本运行
            run.text = self.process_title(text)   # 设置文本为目录项的名称
            run.font.size = Pt(24)  # 设置字体大小为24磅
            run.font.bold = True  # 设置字体为粗体
            run.font.color.rgb = RGBColor(68, 84, 106)  # 设置字体颜色为黑色

            # 添加编号框（左侧显示编号）
            num_box = slide.shapes.add_textbox(left - Inches(1), top, Inches(0.8), height)  # 设置编号框的位置和大小
            num_text_frame = num_box.text_frame  # 获取编号框的文本框架
            num_text_frame.clear()  # 清空任何已有的文本内容

            # 设置编号框中的垂直对齐方式为居中
            num_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # 设置编号框中的文本段落对齐方式为居中
            p = num_text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER  # 水平居中

            # 设置该段落的编号文本
            run = p.add_run()
            run.text = str(i + 1)  # 编号从1开始
            run.font.size = Pt(28)  # 设置编号字体大小为28磅
            run.font.bold = True  # 设置编号为粗体
            run.font.color.rgb = RGBColor(255, 255, 255)  # 设置编号的字体颜色为白色

            # 设置编号框背景颜色为深色
            num_box.fill.solid()  # 填充编号框颜色
            num_box.fill.fore_color.rgb = RGBColor(54, 69, 79)  # 设置背景色为深色（灰蓝色）

            # 设置编号框的边框颜色为透明
            num_box.line.fill.background()  # 设置编号框边框为透明

    # 添加大标题页
    def add_main_title(self, title_text, title_num):
        slide = self.prs_new.slides.add_slide(self.slide_layout[11])

        content = slide.placeholders[0]  # 获取占位符
        content.text = ""  # 先清空原有文本

        # 获取 TextFrame
        text_frame = content.text_frame
        text_frame.clear()  # 清空可能存在的段落

        # 添加段落并设置文本
        p = text_frame.paragraphs[0]
        p.text = title_text
        #p.font.size = Pt(44)  # 调整字体大小
        p.font.bold = True  # 加粗
        p.font.color.rgb = RGBColor(192, 0, 0)  # 设为红色


        # 添加大号背景数字
        bg_text = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(10), Inches(3))
        bg_p = bg_text.text_frame.add_paragraph()
        bg_p.text = title_num
        bg_p.font.size = Pt(275)
        bg_p.font.bold = True
        bg_p.font.color.rgb = RGBColor(255, 192, 192)  # 设为淡红色
        bg_p.alignment = PP_ALIGN.LEFT

    # 添加文本页
    def add_all_text(self, title, text):
        slide = self.prs_new.slides.add_slide(self.slide_layout[8])
        text_len = len(text)
        text_size = 25
        row_text = 34
        if text_len >= 200 and text_len < 300:
            text_size = int(300 / text_len * 25 * 0.8)
            row_text = int(text_len / 400 * 34 * 1.2)
            print(f"文本字体大小：{text_size}")

        elif text_len >= 300 and text_len < 600:
            text_size = int(300 / text_len * 25 * 1)
            row_text = int(text_len / 400 * 34 * 1.2)
            print(f"文本字体大小：{text_size}")

        text_height = 0.5 * (text_len // row_text + 1)
        text_top = 1 + 0.5 * (6 - text_height)

        slide.placeholders[0].text = title
        slide.placeholders[10].top = Inches(text_top)
        slide.placeholders[10].left = Inches(0.5)
        slide.placeholders[10].width = Inches(12.3)
        slide.placeholders[10].height = Inches(text_height)
        slide.placeholders[10].text = text

        for paragraph in slide.placeholders[10].text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(text_size)

    # 添加文本和图片页（文本左-图片右）
    def add_text_image(self, title, text, image_path):
        text_len = len(text)
        print(f"文本长度：{text_len}")
        # 创建新的幻灯片
        slide = self.prs_new.slides.add_slide(self.slide_layout[8])

        # 设置标题
        slide.placeholders[0].text = title

        # 设置文本的占位符
        text_placeholder = slide.placeholders[10]
        text_placeholder.text = text

        row_text = 15
        text_size = 25
        if text_len >= 200 and text_len < 400:
            text_size = int(200 / text_len * 25 * 1.1)
            print(f"文本字体大小：{text_size}")
            row_text = int(text_len / 200 * 15 * 1.2)
            print(f"文本每行字数：{row_text}")
        elif text_len >= 400 and text_len < 600:
            text_size = int(200 / text_len * 25 * 1.5)
            print(f"文本字体大小：{text_size}")
            row_text = int(text_len / 200 * 15 * 1.2)
            print(f"文本每行字数：{row_text}")

        for paragraph in text_placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(text_size)

        # 设置文本框的位置和大小
        text_height = 0.5 * (text_len // row_text + 1)
        text_top = 1 + 0.5 * (6 - text_height)
        print(f"文本框高度：{text_height}，文本框顶部距离：{text_top}")
        text_placeholder.left = Inches(0.5)
        text_placeholder.top = Inches(text_top)
        text_placeholder.width = Inches(6)
        text_placeholder.height = Inches(text_height)

        # 获取图片的大小
        with Image.open(image_path) as img:
            img_width, img_height = img.size

        # 计算PPT页面的宽度和最大图片宽度（PPT整体宽度的一半）
        img_width_alter = 600

        # 按照宽度调整图片比例

        scale_factor = img_width_alter / img_width
        img_width = img_width_alter
        img_height = int(img_height * scale_factor)  # 按比例调整高度

        # 计算图片的位置：使其位于PPT的右侧

        p_height =  (8 - img_height / 96) / 2
        p_width = 6.5  + (6.5 - img_width / 96) / 2
        left = Inches(p_width)  # 图片距离左侧的距离
        top = Inches(p_height)  # 图片距离顶部的距离

        # 插入图片
        slide.shapes.add_picture(image_path, left, top, width=Inches(img_width / 96), height=Inches(img_height / 96))

    # 添加文本和图片页（图片*2）
    def add_text_double_image(self, title, text, image_path_1, image_path_2):
        text_len = len(text)
        print(f"文本长度：{text_len}")
        # 创建新的幻灯片
        slide = self.prs_new.slides.add_slide(self.slide_layout[8])

        # 设置标题
        slide.placeholders[0].text = title

        # 设置文本的占位符
        text_placeholder = slide.placeholders[10]
        text_placeholder.text = text

        row_text = 15
        text_size = 25
        if text_len >= 100 and text_len < 200:
            text_size = int(100 / text_len * 25 * 1.1)
            print(f"文本字体大小：{text_size}")
            row_text = int(text_len / 100 * 15 * 1.2)
            print(f"文本每行字数：{row_text}")
        elif text_len >= 200 and text_len < 300:
            text_size = int(100 / text_len * 25 * 1.5)
            print(f"文本字体大小：{text_size}")
            row_text = int(text_len / 100 * 15 * 1.2)
            print(f"文本每行字数：{row_text}")


        for paragraph in text_placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(text_size)

        # 设置文本框的位置和大小
        text_height = 0.5 * (text_len // row_text + 1)
        text_top = 3.7 + 0.5 * (4 - text_height)
        print(f"文本框高度：{text_height}，文本框顶部距离：{text_top}")
        text_placeholder.left = Inches(0.3)
        text_placeholder.top = Inches(text_top)
        text_placeholder.width = Inches(6)
        text_placeholder.height = Inches(text_height)


        # 获取图片的大小
        with Image.open(image_path_1) as img_1:
            img_1_width, img_1_height = img_1.size

        with Image.open(image_path_2) as img_2:
            img_2_width, img_2_height = img_2.size

        # 更扁平的图像放在左侧文本框上方，更瘦长的放在右侧
        if img_1_width / img_1_height > img_2_width / img_2_height:
            image_left_path = image_path_1
            image_right_path = image_path_2
            img_left_width = img_1_width
            img_left_height = img_1_height
            img_right_width = img_2_width
            img_right_height = img_2_height
        else:
            image_left_path = image_path_2
            image_right_path = image_path_1
            img_left_width = img_2_width
            img_left_height = img_2_height
            img_right_width = img_1_width
            img_right_height = img_1_height


        # 计算PPT页面的宽度和最大图片宽度（PPT整体宽度的一半）
        img_left_width_alter = 600

        # 按照宽度调整图片比例

        scale_factor_left = img_left_width_alter / img_left_width
        img_left_width = img_left_width_alter
        img_left_height = int(img_left_height * scale_factor_left)  # 按比例调整高度

        # 计算图片的位置：使img_left位于PPT的左侧，文本框上方

        p_left_height =  1 + (3 - img_left_height / 96) / 2
        p_left_width = (6.5 - img_left_width / 96) / 2
        left = Inches(p_left_width)  # 图片距离左侧的距离
        top = Inches(p_left_height)  # 图片距离顶部的距离

        # 插入图片
        slide.shapes.add_picture(image_left_path, left, top, width=Inches(img_left_width / 96), height=Inches(img_left_height / 96))

        # 计算PPT页面的宽度和最大图片宽度（PPT整体宽度的一半）
        img_right_width_alter = 600

        # 按照宽度调整图片比例

        scale_factor_right = img_right_width_alter / img_right_width
        img_right_width = img_right_width_alter
        img_right_height = int(img_right_height * scale_factor_right)

        # 计算图片的位置：使img_right位于PPT的右侧，文本框上方

        p_right_height =  (8 - img_right_height / 96) / 2
        p_right_width = 6.5 + (6.5 - img_right_width / 96) / 2
        left = Inches(p_right_width)
        top = Inches(p_right_height)

        # 插入图片
        slide.shapes.add_picture(image_right_path, left, top, width=Inches(img_right_width / 96), height=Inches(img_right_height / 96))

    # 添加图片页
    def add_all_image(self, title, image_path):

        slide = self.prs_new.slides.add_slide(self.slide_layout[8])

        slide.placeholders[0].text = title

        text_placeholder = slide.placeholders[10]
        text_placeholder.element.getparent().remove(text_placeholder.element)
        # 获取图片的大小
        with Image.open(image_path) as img:
            img_width, img_height = img.size

        # 计算PPT页面的宽度和最大图片宽度（PPT整体宽度的一半）
        img_width_alter = 1100
        img_height_alter = 550

        # 调整图片大小
        if img_width / img_height > img_width_alter / img_height_alter:
            # 如果图片的宽高比大于PPT页面的宽高比，则按照宽度调整图片比例
            scale_factor = img_width_alter / img_width
            img_width = img_width_alter
            img_height = int(img_height * scale_factor)
        else:
            # 如果图片的宽高比小于PPT页面的宽高比，则按照高度调整图片比例
            scale_factor = img_height_alter / img_height
            img_height = img_height_alter
            img_width = int(img_width * scale_factor)


        # 计算图片的位置：使其位于PPT的中央
        p_height = (8 - img_height / 96) / 2
        p_width =  (13 - img_width / 96) / 2
        left = Inches(p_width)  # 图片距离左侧的距离
        top = Inches(p_height)  # 图片距离顶部的距离

        # 插入图片
        slide.shapes.add_picture(image_path, left, top, width=Inches(img_width / 96), height=Inches(img_height / 96))

    # 添加双图片页
    def add_double_image(self, title, image_path_1, image_path_2):

        slide = self.prs_new.slides.add_slide(self.slide_layout[8])

        slide.placeholders[0].text = title

        text_placeholder = slide.placeholders[10]
        text_placeholder.element.getparent().remove(text_placeholder.element)
        # 获取图片的大小
        with Image.open(image_path_1) as img_1:
            img_1_width, img_1_height = img_1.size

        # 计算PPT页面的宽度和最大图片宽度（PPT整体宽度的一半）
        img_width_alter = 550
        img_height_alter = 550

        # 调整图片大小
        if img_1_width / img_1_height > img_width_alter / img_height_alter:
            # 如果图片的宽高比大于PPT页面的宽高比，则按照宽度调整图片比例
            scale_factor = img_width_alter / img_1_width
            img_1_width = img_width_alter
            img_1_height = int(img_1_height * scale_factor)
        else:
            # 如果图片的宽高比小于PPT页面的宽高比，则按照高度调整图片比例
            scale_factor = img_height_alter / img_1_height
            img_1_height = img_height_alter
            img_1_width = int(img_1_width * scale_factor)


        # 计算图片的位置：使其位于PPT的中央
        p_1_height = (8 - img_1_height / 96) / 2
        p_1_width = 0.3 + (6 - img_1_width / 96) / 2
        left = Inches(p_1_width)  # 图片距离左侧的距离
        top = Inches(p_1_height)  # 图片距离顶部的距离

        # 插入图片
        slide.shapes.add_picture(image_path_1, left, top, width=Inches(img_1_width / 96), height=Inches(img_1_height / 96))

        with Image.open(image_path_2) as img_2:
            img_2_width, img_2_height = img_2.size

        # 调整图片大小
        if img_2_width / img_2_height > img_width_alter / img_height_alter:
            # 如果图片的宽高比大于PPT页面的宽高比，则按照宽度调整图片比例
            scale_factor = img_width_alter / img_2_width
            img_2_width = img_width_alter
            img_2_height = int(img_2_height * scale_factor)
        else:
            # 如果图片的宽高比小于PPT页面的宽高比，则按照高度调整图片比例
            scale_factor = img_height_alter / img_2_height
            img_2_height = img_height_alter
            print("add double image")
            img_2_width = int(img_2_width * scale_factor)


        # 计算图片的位置：使其位于PPT的中央
        p_2_height = (8 - img_2_height / 96) / 2
        p_2_width = 7 + (6 - img_2_width / 96) / 2
        left = Inches(p_2_width)  # 图片距离左侧的距离
        top = Inches(p_2_height)  # 图片距离顶部的距离

        # 插入图片
        slide.shapes.add_picture(image_path_2, left, top, width=Inches(img_2_width / 96), height=Inches(img_2_height / 96))

    # 添加感谢页
    def add_thanks(self, text="感谢您的观看"):
        slide = self.prs_new.slides.add_slide(self.slide_layout[13])

        slide.placeholders[0].text = text



    # 保存PPT
    def save_ppt(self, output_path):
        self.prs_new.save(output_path)
        print(f"PPT 已成功保存到 {output_path}")